"""
Create PowerPoint from slide screenshots
"""
from pptx import Presentation
from pptx.util import Inches
import os

# Create widescreen presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide images - update these paths
slide_images = [
    r"C:\Users\sanja\.gemini\antigravity\brain\c83413a7-d5ae-42c0-8aeb-9e959ce0c9a2\uploaded_image_0_1767691457225.png",
    r"C:\Users\sanja\.gemini\antigravity\brain\c83413a7-d5ae-42c0-8aeb-9e959ce0c9a2\uploaded_image_1_1767691457225.png",
    r"C:\Users\sanja\.gemini\antigravity\brain\c83413a7-d5ae-42c0-8aeb-9e959ce0c9a2\uploaded_image_2_1767691457225.png",
    r"C:\Users\sanja\.gemini\antigravity\brain\c83413a7-d5ae-42c0-8aeb-9e959ce0c9a2\uploaded_image_3_1767691457225.png",
    r"C:\Users\sanja\.gemini\antigravity\brain\c83413a7-d5ae-42c0-8aeb-9e959ce0c9a2\uploaded_image_4_1767691457225.png",
]

# Add each image as a full-slide
for img_path in slide_images:
    if os.path.exists(img_path):
        blank_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(blank_layout)
        
        # Add image to fill entire slide
        slide.shapes.add_picture(
            img_path,
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
        print(f"✅ Added: {os.path.basename(img_path)}")
    else:
        print(f"❌ Not found: {img_path}")

# Save
output_path = r"C:\Users\sanja\OneDrive\Desktop\pcds-core\PCDS_PITCH_DECK.pptx"
prs.save(output_path)
print(f"\n🎉 PowerPoint saved: {output_path}")
