"""
Generate PCDS Pitch Deck PowerPoint
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RgbColor(1, 2, 5)
CYAN = RgbColor(0, 242, 255)
BLUE = RgbColor(0, 120, 212)
RED = RgbColor(255, 45, 85)
WHITE = RgbColor(255, 255, 255)
GRAY = RgbColor(128, 128, 128)

def add_dark_slide(prs):
    """Add a blank slide with dark background"""
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)
    
    # Add dark rectangle as background
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    
    # Send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """Add a text box to slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

# SLIDE 1: Problem
slide1 = add_dark_slide(prs)
add_text_box(slide1, 0.8, 0.5, 3, 0.5, "THE CRISIS", 14, True, RED)
add_text_box(slide1, 0.8, 1.0, 6, 2, "Every 39\nSeconds", 72, True, WHITE)
add_text_box(slide1, 0.8, 3.5, 5, 1.5, "A cyberattack happens. SMBs are the \"soft underbelly\" of the economy—too small for enterprise security, too vital to fail.", 18, False, GRAY)
add_text_box(slide1, 0.8, 5.0, 2, 1, "60%", 48, True, RED)
add_text_box(slide1, 2.0, 5.2, 4, 0.8, "of SMBs close within 6 months of a breach", 14, False, GRAY)
add_text_box(slide1, 0.8, 6.0, 2, 1, "$50K+", 48, True, RgbColor(255, 165, 0))
add_text_box(slide1, 2.2, 6.2, 4, 0.8, "per year for enterprise security. 99% can't afford it.", 14, False, GRAY)
add_text_box(slide1, 12.5, 0.3, 1, 0.3, "01 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# SLIDE 2: Solution
slide2 = add_dark_slide(prs)
add_text_box(slide2, 0.8, 0.5, 5, 0.4, "☁️ POWERED BY MICROSOFT AZURE", 12, True, BLUE)
add_text_box(slide2, 0.8, 1.0, 5, 1, "PCDS", 72, True, CYAN)
add_text_box(slide2, 0.8, 2.2, 5, 0.5, "Predictive Cyber Defence System for SMBs", 20, False, GRAY)
add_text_box(slide2, 0.8, 3.2, 6, 0.5, "🧠 5-Model ML Ensemble", 20, True, WHITE)
add_text_box(slide2, 0.8, 3.7, 6, 0.4, "Trained on 5.3M+ attack samples. 88.3% accuracy.", 14, False, GRAY)
add_text_box(slide2, 0.8, 4.3, 6, 0.5, "🤖 AI Security Copilot", 20, True, WHITE)
add_text_box(slide2, 0.8, 4.8, 6, 0.4, "Azure OpenAI explains threats in plain English.", 14, False, GRAY)
add_text_box(slide2, 0.8, 5.4, 6, 0.5, "⚡ Automated Response", 20, True, WHITE)
add_text_box(slide2, 0.8, 5.9, 6, 0.4, "<2ms latency. Instant isolation. No human delay.", 14, False, GRAY)
add_text_box(slide2, 8, 3.5, 2, 1, "88.3%", 36, True, CYAN, PP_ALIGN.CENTER)
add_text_box(slide2, 8, 4.2, 2, 0.3, "ACCURACY", 10, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide2, 10, 3.5, 2, 1, "<2ms", 36, True, BLUE, PP_ALIGN.CENTER)
add_text_box(slide2, 10, 4.2, 2, 0.3, "LATENCY", 10, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide2, 8, 5.0, 2, 1, "<3%", 36, True, RgbColor(0, 255, 0), PP_ALIGN.CENTER)
add_text_box(slide2, 8, 5.7, 2, 0.3, "FALSE POSITIVE", 10, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide2, 10, 5.0, 2, 1, "5.3M+", 36, True, RgbColor(160, 32, 240), PP_ALIGN.CENTER)
add_text_box(slide2, 10, 5.7, 2, 0.3, "SAMPLES", 10, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide2, 12.5, 0.3, 1, 0.3, "02 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# SLIDE 3: Architecture
slide3 = add_dark_slide(prs)
add_text_box(slide3, 0.5, 0.5, 12, 0.4, "TECHNICAL ARCHITECTURE", 12, True, BLUE, PP_ALIGN.CENTER)
add_text_box(slide3, 0.5, 1.0, 12, 0.8, "Built on Microsoft Azure", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide3, 1, 2.2, 2.5, 0.4, "📡 Data Ingestion", 16, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide3, 1, 2.6, 2.5, 0.3, "Azure Monitor", 12, False, BLUE, PP_ALIGN.CENTER)
add_text_box(slide3, 3.7, 2.4, 0.5, 0.4, "→", 24, False, CYAN, PP_ALIGN.CENTER)
add_text_box(slide3, 4.2, 2.2, 2.5, 0.4, "🧠 ML Detection", 16, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide3, 4.2, 2.6, 2.5, 0.3, "Azure ML + GPU", 12, False, CYAN, PP_ALIGN.CENTER)
add_text_box(slide3, 6.9, 2.4, 0.5, 0.4, "→", 24, False, CYAN, PP_ALIGN.CENTER)
add_text_box(slide3, 7.4, 2.2, 2.5, 0.4, "⚡ Risk Scoring", 16, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide3, 7.4, 2.6, 2.5, 0.3, "88.3% Accuracy", 12, False, RgbColor(160, 32, 240), PP_ALIGN.CENTER)
add_text_box(slide3, 10.1, 2.4, 0.5, 0.4, "→", 24, False, CYAN, PP_ALIGN.CENTER)
add_text_box(slide3, 10.6, 2.2, 2.5, 0.4, "🛡️ Auto Response", 16, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide3, 10.6, 2.6, 2.5, 0.3, "Azure Functions", 12, False, RgbColor(0, 255, 0), PP_ALIGN.CENTER)
add_text_box(slide3, 0.5, 3.5, 12, 0.4, "AZURE SERVICES INTEGRATED", 10, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide3, 0.5, 4.0, 12, 0.5, "Azure OpenAI  •  Azure Functions  •  Azure Kubernetes Service  •  Azure Blob Storage  •  Azure Monitor  •  Azure SQL", 14, False, BLUE, PP_ALIGN.CENTER)
add_text_box(slide3, 12.5, 0.3, 1, 0.3, "03 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# SLIDE 4: Features
slide4 = add_dark_slide(prs)
add_text_box(slide4, 0.5, 0.5, 12, 0.4, "WHAT USERS SEE", 12, True, CYAN, PP_ALIGN.CENTER)
add_text_box(slide4, 0.5, 1.0, 12, 0.8, "Built for Non-Technical Owners", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide4, 1, 2.2, 5.5, 0.5, "📊 Smart Dashboard", 20, True, WHITE)
add_text_box(slide4, 1, 2.7, 5.5, 0.5, "Real-time risk score, active threats, simple status colors.", 14, False, GRAY)
add_text_box(slide4, 7, 2.2, 5.5, 0.5, "🤖 AI Security Copilot", 20, True, WHITE)
add_text_box(slide4, 7, 2.7, 5.5, 0.5, "Plain-language answers via Azure OpenAI.", 14, False, GRAY)
add_text_box(slide4, 1, 3.8, 5.5, 0.5, "⚡ Automated Playbooks", 20, True, WHITE)
add_text_box(slide4, 1, 4.3, 5.5, 0.5, "Instant isolation, IP blocking, team alerts.", 14, False, GRAY)
add_text_box(slide4, 7, 3.8, 5.5, 0.5, "🎯 MITRE ATT&CK Mapping", 20, True, WHITE)
add_text_box(slide4, 7, 4.3, 5.5, 0.5, "Every threat mapped to industry framework.", 14, False, GRAY)
add_text_box(slide4, 12.5, 0.3, 1, 0.3, "04 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# SLIDE 5: Validation
slide5 = add_dark_slide(prs)
add_text_box(slide5, 0.5, 0.5, 12, 0.4, "EARLY VALIDATION", 12, True, RgbColor(0, 255, 0), PP_ALIGN.CENTER)
add_text_box(slide5, 0.5, 1.0, 12, 0.8, "Real Feedback from India", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide5, 1, 2.2, 11, 0.6, "🏦 \"We surveyed SMB banks across India — they loved the simplicity and are ready to adopt PCDS.\"", 16, False, WHITE)
add_text_box(slide5, 1.3, 2.8, 5, 0.3, "SMB Bank Survey • Strong Purchase Intent", 11, False, GRAY)
add_text_box(slide5, 1, 3.4, 11, 0.6, "💼 \"Enterprise-level protection that fits our budget. We want this product.\"", 16, False, WHITE)
add_text_box(slide5, 1.3, 4.0, 5, 0.3, "Cooperative Bank Manager • Rural India", 11, False, GRAY)
add_text_box(slide5, 1, 4.6, 11, 0.6, "🎓 \"Educational institutions loved how easy it is to protect student data.\"", 16, False, WHITE)
add_text_box(slide5, 1.3, 5.2, 5, 0.3, "Education Sector Survey • Schools & Colleges", 11, False, GRAY)
add_text_box(slide5, 1, 5.8, 11, 0.6, "⭐ \"Stunning work! We will deploy this in our college soon.\"", 16, False, WHITE)
add_text_box(slide5, 1.3, 6.4, 5, 0.3, "Sandeep Babu Challa • Industry Mentor • First Deployment", 11, False, BLUE)
add_text_box(slide5, 12.5, 0.3, 1, 0.3, "05 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# SLIDE 6: Competitive
slide6 = add_dark_slide(prs)
add_text_box(slide6, 0.5, 0.5, 12, 0.4, "MARKET GAP", 12, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 0.5, 1.0, 12, 0.8, "Enterprise Security is a Luxury", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide6, 1.5, 2.5, 2.5, 0.4, "DARKTRACE", 12, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 1.5, 3.0, 2.5, 0.8, "$100k+", 36, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 1.5, 3.8, 2.5, 0.3, "Annual", 10, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 4.5, 2.5, 2.5, 0.4, "VECTRA AI", 12, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 4.5, 3.0, 2.5, 0.8, "$80k+", 36, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 4.5, 3.8, 2.5, 0.3, "Enterprise", 10, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 7.5, 2.5, 2.5, 0.4, "CROWDSTRIKE", 12, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 7.5, 3.0, 2.5, 0.8, "$50k+", 36, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 7.5, 3.8, 2.5, 0.3, "Core Platform", 10, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 10.5, 2.5, 2.5, 0.4, "🛡️ PCDS", 14, True, CYAN, PP_ALIGN.CENTER)
add_text_box(slide6, 10.5, 3.0, 2.5, 0.8, "$1,200", 36, True, CYAN, PP_ALIGN.CENTER)
add_text_box(slide6, 10.5, 3.8, 2.5, 0.3, "Per Year", 10, False, CYAN, PP_ALIGN.CENTER)
add_text_box(slide6, 0.5, 5.0, 12, 0.6, "\"Democratizing elite cybersecurity for the 99%.\"", 24, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide6, 12.5, 0.3, 1, 0.3, "06 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# SLIDE 7: GTM
slide7 = add_dark_slide(prs)
add_text_box(slide7, 0.5, 0.5, 12, 0.4, "BUSINESS MODEL", 12, True, RgbColor(160, 32, 240), PP_ALIGN.CENTER)
add_text_box(slide7, 0.5, 1.0, 12, 0.8, "Go-to-Market Strategy", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide7, 1, 2.2, 5, 0.5, "Distribution Channels", 20, True, WHITE)
add_text_box(slide7, 1, 2.8, 5, 0.4, "☁️ Azure Marketplace", 16, True, CYAN)
add_text_box(slide7, 1, 3.2, 5, 0.4, "Direct to SMBs using Microsoft 365", 12, False, GRAY)
add_text_box(slide7, 1, 3.8, 5, 0.4, "🤝 MSP Partnerships", 16, True, RgbColor(160, 32, 240))
add_text_box(slide7, 1, 4.2, 5, 0.4, "Managed Service Providers resell to SMBs", 12, False, GRAY)
add_text_box(slide7, 7, 2.2, 5, 0.5, "Pricing Model", 20, True, WHITE)
add_text_box(slide7, 7, 2.8, 2.5, 0.4, "Starter: $99/mo", 16, True, CYAN)
add_text_box(slide7, 9.5, 2.8, 2.5, 0.3, "1-25 endpoints", 11, False, GRAY)
add_text_box(slide7, 7, 3.4, 2.5, 0.4, "Growth: $299/mo", 16, True, BLUE)
add_text_box(slide7, 9.5, 3.4, 2.5, 0.3, "26-100 endpoints", 11, False, GRAY)
add_text_box(slide7, 7, 4.0, 2.5, 0.4, "Enterprise: Custom", 16, True, RgbColor(160, 32, 240))
add_text_box(slide7, 9.5, 4.0, 2.5, 0.3, "100+ endpoints", 11, False, GRAY)
add_text_box(slide7, 0.5, 5.2, 12, 0.3, "LAUNCH TIMELINE", 10, True, GRAY, PP_ALIGN.CENTER)
add_text_box(slide7, 1.5, 5.6, 3, 0.4, "Q1 2026", 16, True, CYAN, PP_ALIGN.CENTER)
add_text_box(slide7, 1.5, 6.0, 3, 0.3, "Final pilots + Marketplace", 11, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide7, 5, 5.6, 3, 0.4, "Q2 2026", 16, True, BLUE, PP_ALIGN.CENTER)
add_text_box(slide7, 5, 6.0, 3, 0.3, "Beta with 10 MSPs", 11, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide7, 8.5, 5.6, 3, 0.4, "Q3 2026", 16, True, RgbColor(160, 32, 240), PP_ALIGN.CENTER)
add_text_box(slide7, 8.5, 6.0, 3, 0.3, "Scale: India + SE Asia", 11, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide7, 12.5, 0.3, 1, 0.3, "07 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# SLIDE 8: Why Now
slide8 = add_dark_slide(prs)
add_text_box(slide8, 0.5, 0.8, 12, 1, "Why PCDS Wins", 60, True, CYAN, PP_ALIGN.CENTER)
add_text_box(slide8, 1, 2.2, 5.5, 0.5, "⏰ Timing", 18, True, WHITE)
add_text_box(slide8, 1, 2.7, 5.5, 0.4, "Post-2024 attacks have SMBs desperate for solutions", 12, False, GRAY)
add_text_box(slide8, 7, 2.2, 5.5, 0.5, "☁️ Azure Advantage", 18, True, WHITE)
add_text_box(slide8, 7, 2.7, 5.5, 0.4, "Native integration. Marketplace visibility.", 12, False, GRAY)
add_text_box(slide8, 1, 3.5, 5.5, 0.5, "👥 Team", 18, True, WHITE)
add_text_box(slide8, 1, 4.0, 5.5, 0.4, "CS students with cybersecurity passion", 12, False, GRAY)
add_text_box(slide8, 7, 3.5, 5.5, 0.5, "📈 Market", 18, True, WHITE)
add_text_box(slide8, 7, 4.0, 5.5, 0.4, "2M+ SMBs in India. Global TAM: $50B+", 12, False, GRAY)
add_text_box(slide8, 0.5, 4.8, 12, 0.5, "With Imagine Cup Support", 20, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide8, 0.5, 5.3, 12, 0.5, "✓ Accelerate Azure Marketplace  •  ✓ Build MSP partnerships  •  ✓ Protect thousands of SMBs", 14, False, CYAN, PP_ALIGN.CENTER)
add_text_box(slide8, 0.5, 6.0, 12, 0.6, "Let's build enterprise security for everyone.", 28, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide8, 0.5, 6.8, 12, 0.4, "Team SURAKSHA AI  •  Imagine Cup 2026", 14, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide8, 12.5, 0.3, 1, 0.3, "08 / 08", 12, False, GRAY, PP_ALIGN.RIGHT)

# Save
output_path = "PCDS_PITCH_DECK.pptx"
prs.save(output_path)
print(f"✅ PowerPoint saved: {output_path}")
print(f"   Location: {os.path.abspath(output_path)}")
